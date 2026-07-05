# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Generate the ProRex Consultancy branded PowerPoint template into brand/.

Run with:  uv run scripts/generate_brand_deck.py

The output is NOT committed (brand/ is gitignored) — it is encrypted into
public/vault/ by scripts/vault_encrypt.py.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "brand"
OUT.mkdir(exist_ok=True)

NAVY = RGBColor(0x0C, 0x2D, 0x4D)
NAVY_DEEP = RGBColor(0x09, 0x22, 0x3A)
INK = RGBColor(0x16, 0x32, 0x4F)
ICE = RGBColor(0xD8, 0xE9, 0xF9)
DIM = RGBColor(0x8F, 0xB3, 0xD1)
ACCENT = RGBColor(0x6F, 0xD3, 0xFF)
ACCENT_LIGHT = RGBColor(0x0D, 0x6F, 0xB8)
PAPER = RGBColor(0xF7, 0xF5, 0xEE)
LINE_LIGHT = RGBColor(0xC4, 0xCF, 0xDA)
STAMP = RGBColor(0xFF, 0x6B, 0x6B)
STAMP_LIGHT = RGBColor(0xD4, 0x3D, 0x3D)

MONO = "Courier New"
SANS = "Arial"

W, H = Inches(13.333), Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    return sh


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, font, size, color, bold, spacing)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, font, size, color, bold, spacing in para:
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            if spacing:
                # letter spacing in 1/100 pt
                r.font._rPr.set("spc", str(spacing))
    return tb


def grid(slide, color, alpha_step=Inches(0.5), weight=0.25):
    """Blueprint grid of thin lines across the slide."""
    from pptx.enum.shapes import MSO_CONNECTOR

    x = Emu(0)
    while x < W:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Emu(0), x, H)
        ln.line.color.rgb = color
        ln.line.width = Pt(weight)
        ln.shadow.inherit = False
        x = Emu(int(x) + int(alpha_step))
    y = Emu(0)
    while y < H:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(0), y, W, y)
        ln.line.color.rgb = color
        ln.line.width = Pt(weight)
        ln.shadow.inherit = False
        y = Emu(int(y) + int(alpha_step))


def frame(slide, ink, sheet_no, title_label, dark=True):
    """Drawing-sheet frame with header strip and title block footer."""
    m = Inches(0.25)
    rect(slide, m, m, W - 2 * m, H - 2 * m, fill=None, line=ink, line_w=1.6)
    # header strip
    rect(slide, m, m, W - 2 * m, Inches(0.42), fill=None, line=ink, line_w=1.0)
    dimc = DIM if dark else RGBColor(0x7A, 0x8E, 0xA3)
    acc = ACCENT if dark else ACCENT_LIGHT
    text(slide, m + Inches(0.2), m + Inches(0.08), Inches(4), Inches(0.3),
         [[("DOC NO. PRX-____ · REV ____", MONO, 10, dimc, False, 100)]])
    text(slide, W / 2 - Inches(3), m + Inches(0.08), Inches(6), Inches(0.3),
         [[(title_label, MONO, 10, acc, True, 200)]], align=PP_ALIGN.CENTER)
    text(slide, W - m - Inches(2.6), m + Inches(0.08), Inches(2.4), Inches(0.3),
         [[(sheet_no, MONO, 10, dimc, False, 150)]], align=PP_ALIGN.RIGHT)
    # title block footer
    fy = H - m - Inches(0.42)
    rect(slide, m, fy, W - 2 * m, Inches(0.42), fill=None, line=ink, line_w=1.0)
    text(slide, m + Inches(0.2), fy + Inches(0.08), Inches(6), Inches(0.3),
         [[("PROREX CONSULTANCY — DATA & PLATFORM ENGINEERING", MONO, 9, dimc, False, 100)]])
    text(slide, W - m - Inches(5.2), fy + Inches(0.08), Inches(5), Inches(0.3),
         [[("MEASURE TWICE · BUILD ONCE", MONO, 9, acc, True, 150)]], align=PP_ALIGN.RIGHT)


def main() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---------------- 1. Title (dark)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=NAVY)
    grid(s, NAVY_DEEP)
    frame(s, ICE, "SHT 01", "PRESENTATION — TITLE SHEET")
    text(s, Inches(1.0), Inches(2.1), Inches(9), Inches(0.4),
         [[("FIG. 001 — WORKING SESSION", MONO, 13, ACCENT, True, 300)]])
    text(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.9),
         [[("PRESENTATION TITLE", SANS, 54, ICE, True, 100)],
          [("Subtitle or client name goes here", SANS, 20, DIM, False, 0)]], line_spacing=1.05)
    text(s, Inches(1.0), Inches(5.3), Inches(9), Inches(0.8),
         [[("MISJA PRONK · PROREX CONSULTANCY", MONO, 12, ICE, True, 150)],
          [("DATE: __________ · AUDIENCE: __________", MONO, 11, DIM, False, 100)]], line_spacing=1.5)

    # ---------------- 2. Agenda (light)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=PAPER)
    grid(s, RGBColor(0xE9, 0xE5, 0xD8))
    frame(s, INK, "SHT 02", "DRAWING SET — CONTENTS", dark=False)
    text(s, Inches(1.0), Inches(1.0), Inches(10), Inches(0.7),
         [[("Agenda", SANS, 36, INK, True, 0)]])
    items = ["Context & goal", "Current situation", "Proposed approach", "Architecture", "Planning & next steps"]
    for i, item in enumerate(items):
        y = Inches(2.0 + i * 0.85)
        rect(s, Inches(1.0), y, Inches(0.55), Inches(0.55), fill=None, line=ACCENT_LIGHT, line_w=1.2)
        text(s, Inches(1.0), y + Inches(0.12), Inches(0.55), Inches(0.35),
             [[(f"{i+1:02d}", MONO, 14, STAMP_LIGHT, True, 0)]], align=PP_ALIGN.CENTER)
        text(s, Inches(1.85), y + Inches(0.08), Inches(9), Inches(0.45),
             [[(item, SANS, 18, INK, False, 0)]])

    # ---------------- 3. Section divider (dark)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=NAVY)
    grid(s, NAVY_DEEP)
    frame(s, ICE, "SHT 03", "SECTION DIVIDER")
    text(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.6),
         [[("01 ", MONO, 40, STAMP, True, 0), ("SECTION TITLE", SANS, 44, ICE, True, 100)],
          [("One line on what this section covers", SANS, 18, DIM, False, 0)]], line_spacing=1.2)

    # ---------------- 4. Content: two-column (light)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=PAPER)
    grid(s, RGBColor(0xE9, 0xE5, 0xD8))
    frame(s, INK, "SHT 04", "DETAIL — CONTENT")
    text(s, Inches(1.0), Inches(1.0), Inches(10), Inches(0.7),
         [[("Content slide title", SANS, 36, INK, True, 0)]])
    text(s, Inches(1.0), Inches(2.1), Inches(5.6), Inches(4.0),
         [[("The point in one sentence.", SANS, 18, ACCENT_LIGHT, True, 0)],
          [("", SANS, 8, INK, False, 0)],
          [("Supporting paragraph: keep it short and concrete. Replace this text with the actual argument, evidence or story.", SANS, 15, INK, False, 0)],
          [("", SANS, 8, INK, False, 0)],
          [("— First supporting point", SANS, 15, INK, False, 0)],
          [("— Second supporting point", SANS, 15, INK, False, 0)],
          [("— Third supporting point", SANS, 15, INK, False, 0)]], line_spacing=1.25)
    # visual placeholder panel
    rect(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(3.9), fill=RGBColor(0xFF, 0xFF, 0xFF), line=INK, line_w=1.2)
    text(s, Inches(7.1), Inches(3.75), Inches(5.2), Inches(0.6),
         [[("[ DIAGRAM / SCREENSHOT / CHART ]", MONO, 12, RGBColor(0x7A, 0x8E, 0xA3), False, 150)]],
         align=PP_ALIGN.CENTER)

    # ---------------- 5. Stats (light)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=PAPER)
    grid(s, RGBColor(0xE9, 0xE5, 0xD8))
    frame(s, INK, "SHT 05", "SCHEDULE OF QUANTITIES")
    text(s, Inches(1.0), Inches(1.0), Inches(10), Inches(0.7),
         [[("The numbers", SANS, 36, INK, True, 0)]])
    stats = [("42%", "METRIC ONE"), ("3×", "METRIC TWO"), ("€ 1.2M", "METRIC THREE")]
    for i, (value, label) in enumerate(stats):
        x = Inches(1.0 + i * 3.9)
        rect(s, x, Inches(2.4), Inches(3.4), Inches(2.2), fill=RGBColor(0xFF, 0xFF, 0xFF), line=LINE_LIGHT, line_w=1.0)
        text(s, x, Inches(2.75), Inches(3.4), Inches(1.0),
             [[(value, SANS, 60, ACCENT_LIGHT, True, 0)]], align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.95), Inches(3.4), Inches(0.4),
             [[(label, MONO, 12, RGBColor(0x7A, 0x8E, 0xA3), True, 200)]], align=PP_ALIGN.CENTER)
    text(s, Inches(1.0), Inches(5.3), Inches(11), Inches(0.5),
         [[("One sentence interpreting the numbers for the audience.", SANS, 15, INK, False, 0)]])

    # ---------------- 6. Architecture (light)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=PAPER)
    grid(s, RGBColor(0xE9, 0xE5, 0xD8))
    frame(s, INK, "SHT 06", "TYPICAL DEPLOYMENT")
    text(s, Inches(1.0), Inches(1.0), Inches(10), Inches(0.7),
         [[("Architecture", SANS, 36, INK, True, 0)]])
    from pptx.enum.shapes import MSO_CONNECTOR
    stages = [("SOURCES", "APIs · DBs · FILES"), ("INGEST", "DATA FACTORY"), ("TRANSFORM", "DATABRICKS · dbt"), ("SERVE", "LAKEHOUSE · BI")]
    for i, (label, sub) in enumerate(stages):
        x = Inches(0.9 + i * 3.05)
        rect(s, x, Inches(3.0), Inches(2.5), Inches(1.15), fill=RGBColor(0xFF, 0xFF, 0xFF), line=INK, line_w=1.4)
        text(s, x, Inches(3.22), Inches(2.5), Inches(0.4),
             [[(label, MONO, 15, INK, True, 200)]], align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.62), Inches(2.5), Inches(0.35),
             [[(sub, MONO, 10, RGBColor(0x7A, 0x8E, 0xA3), False, 100)]], align=PP_ALIGN.CENTER)
        if i < 3:
            ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(2.5), Inches(3.575), x + Inches(3.05), Inches(3.575))
            ln.line.color.rgb = ACCENT_LIGHT
            ln.line.width = Pt(1.8)
            ln.shadow.inherit = False
    text(s, Inches(1.0), Inches(5.0), Inches(11), Inches(0.8),
         [[("Replace with the actual architecture — keep boxes at this size, duplicate rows as needed.", SANS, 14, RGBColor(0x7A, 0x8E, 0xA3), False, 0)]])

    # ---------------- 7. Closing (dark)
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=NAVY)
    grid(s, NAVY_DEEP)
    frame(s, ICE, "SHT 07", "CLOSING — SIGN-OFF")
    text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.2),
         [[("Questions & next steps", SANS, 44, ICE, True, 0)]])
    text(s, Inches(1.0), Inches(4.1), Inches(10), Inches(1.4),
         [[("MISJA PRONK — PROREX CONSULTANCY", MONO, 13, ACCENT, True, 150)],
          [("misja@prorexconsultancy.nl · prorexconsultancy.nl · linkedin.com/in/misja-pronk", MONO, 12, DIM, False, 50)],
          [("KVK 85369624", MONO, 11, DIM, False, 150)]], line_spacing=1.6)

    path = OUT / "prorex-deck-template.pptx"
    prs.save(path)
    print(f"wrote {path}")


if __name__ == "__main__":
    main()
